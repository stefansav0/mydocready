from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

# ✅ Output Name
presentation_name = "DocReady_Styled_Presentation.pptx"

# ✅ Colors
PRIMARY_COLOR = RGBColor(37, 99, 235)  # Blue
ACCENT_COLOR = RGBColor(22, 163, 74)   # Green

# ✅ Create Presentation
prs = Presentation()

# ---------- Function to format title ----------
def format_title(shape, font_size=44, bold=True, color=PRIMARY_COLOR):
    text_frame = shape.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Poppins"  # fallback if not available

# ---------- Title Slide ----------
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "DocReady Company Profile"
subtitle.text = "Empowering Digital Document Services"
format_title(title, font_size=54)

# ---------- Vision Slide ----------
section_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(section_layout)
slide.shapes.title.text = "Our Vision"
format_title(slide.shapes.title)
vision_text = (
    "To make government and personal documents easily accessible for everyone "
    "with the help of modern technology and AI-driven tools."
)
p = slide.placeholders[1].text_frame
p.text = vision_text
p.paragraphs[0].font.size = Pt(20)

# ---------- Optional Image Slide ----------
img_path = "your_image.png"  # Replace with your image path
if os.path.exists(img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    left = Inches(2)
    top = Inches(1)
    height = Inches(3)
    slide.shapes.add_picture(img_path, left, top, height=height)

# ---------- Growth Chart Slide ----------
chart_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title + Content layout
chart_slide.shapes.title.text = "User Growth Chart"
format_title(chart_slide.shapes.title)

# ✅ Chart Data
chart_data = CategoryChartData()
chart_data.categories = ['2023', '2024', '2025']
chart_data.add_series('Users', (1500, 3500, 7000))

# ✅ Add Chart
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = chart_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

# ✅ Style Chart
chart.has_title = True
chart.chart_title.text_frame.text = "Yearly User Growth"
chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

# ---------- Services Slide ----------
slide = prs.slides.add_slide(section_layout)
slide.shapes.title.text = "Our Services"
format_title(slide.shapes.title)
services_text = (
    "✔ Aadhaar Card Services\n"
    "✔ PAN Card Updates\n"
    "✔ EPFO & ESIC Assistance\n"
    "✔ Ayushman Bharat Services\n"
    "✔ AI-Powered Document Editing\n"
    "✔ Photo & Signature Resizing"
)
p = slide.placeholders[1].text_frame
p.text = services_text
p.paragraphs[0].font.size = Pt(18)

# ---------- Legal Disclaimer Slide ----------
legal_slide = prs.slides.add_slide(section_layout)
legal_slide.shapes.title.text = "Legal Disclaimer & Terms"
format_title(legal_slide.shapes.title)
legal_text = (
    "1. DocReady provides assistance for documentation purposes only and is not a government entity.\n"
    "2. All data processed is handled securely, and we comply with data protection norms.\n"
    "3. Users are responsible for providing correct and legal documents.\n"
    "4. Our services do not guarantee government approvals; final decision rests with authorities.\n"
    "5. By using our services, you agree to our Terms & Privacy Policy."
)
p = legal_slide.placeholders[1].text_frame
p.text = legal_text
p.paragraphs[0].font.size = Pt(16)

# ---------- Save Presentation ----------
prs.save(presentation_name)
print(f"✅ Presentation saved as {presentation_name}")
